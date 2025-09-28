import boto3, io, os, time
from core.loader_base import LoaderBase
try:
    import pdfplumber
except Exception:
    pdfplumber=None
try:
    import docx
except Exception:
    docx=None
try:
    import pptx
except Exception:
    pptx=None
try:
    import pytesseract
    from PIL import Image
except Exception:
    pytesseract=None

class S3Loader(LoaderBase):
    def __init__(self,bucket,prefix='',aws_region=None,max_objects=None,sleep_on_rate=1.0,max_workers=4):
        self.bucket=bucket; self.prefix=prefix; self.max_objects=max_objects; self.sleep_on_rate=sleep_on_rate; self.max_workers=max_workers
        self.s3=boto3.client('s3', region_name=aws_region) if aws_region else boto3.client('s3')
    def _is_text(self,key):
        text_ext=['.txt','.md','.py','.json','.csv','.yaml','.yml','.html','.htm']
        _,ext=os.path.splitext(key.lower()); return ext in text_ext
    def _extract_pdf(self,body_bytes):
        if pdfplumber is None: return ''
        try:
            with pdfplumber.open(io.BytesIO(body_bytes)) as pdf:
                pages=[p.extract_text() or '' for p in pdf.pages]
            return '\n\n'.join(pages)
        except Exception as e:
            print('pdf extract error',e); return ''
    def _extract_docx(self,body_bytes):
        if docx is None: return ''
        try:
            doc=docx.Document(io.BytesIO(body_bytes)); paras=[p.text for p in doc.paragraphs]; return '\n\n'.join(paras)
        except Exception as e:
            print('docx extract error',e); return ''
    def _extract_pptx(self,body_bytes):
        if pptx is None: return ''
        try:
            prs=pptx.Presentation(io.BytesIO(body_bytes)); texts=[]
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape,'text'): texts.append(shape.text)
            return '\n\n'.join(texts)
        except Exception as e:
            print('pptx extract error',e); return ''
    def _ocr_image(self,body_bytes):
        if pytesseract is None: return ''
        try:
            img=Image.open(io.BytesIO(body_bytes)); return pytesseract.image_to_string(img)
        except Exception as e:
            print('ocr error',e); return ''
    def load(self):
        paginator=self.s3.get_paginator('list_objects_v2'); out=[]; seen=0
        for page in paginator.paginate(Bucket=self.bucket, Prefix=self.prefix):
            for obj in page.get('Contents',[]):
                key=obj['Key']
                if self.max_objects and seen>=self.max_objects: return out
                try:
                    resp=self.s3.get_object(Bucket=self.bucket, Key=key); body=resp['Body'].read()
                except Exception as e:
                    print('s3 get error',key,e); time.sleep(self.sleep_on_rate); continue
                text=''; lower=key.lower()
                if self._is_text(key):
                    try: text=body.decode('utf-8')
                    except Exception:
                        try: text=body.decode('latin-1')
                        except Exception: text=''
                elif lower.endswith('.pdf'):
                    text=self._extract_pdf(body)
                elif lower.endswith('.docx'):
                    text=self._extract_docx(body)
                elif lower.endswith('.pptx'):
                    text=self._extract_pptx(body)
                elif lower.endswith(('.png','.jpg','.jpeg','.tiff','.bmp')):
                    text=self._ocr_image(body)
                else:
                    continue
                if not text: continue
                out.append({'id':f's3://{self.bucket}/{key}','text':text,'meta':{'source':'s3','key':key}}); seen+=1
        return out
